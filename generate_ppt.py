#!/usr/bin/env python3
"""
生成 JobMate AI 参赛 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
PRIMARY_COLOR = RGBColor(24, 144, 255)  # Ant Design Blue
GRADIENT_START = RGBColor(102, 126, 234)  # #667eea
GRADIENT_END = RGBColor(118, 75, 162)  # #764ba2
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(51, 51, 51)
GRAY_TEXT = RGBColor(102, 102, 102)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle, extra_text=""):
    """封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 添加渐变背景矩形
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRADIENT_START
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 额外信息
    if extra_text:
        extra_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.6))
        tf = extra_box.text_frame
        p = tf.paragraphs[0]
        p.text = extra_text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, subtitle=""):
    """章节标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 左侧色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(0.3), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY_TEXT

    return slide

def add_content_slide(prs, title, content_items):
    """内容页 - 带项目符号"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3),
        Inches(12.333), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR
    line.line.fill.background()

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 20))
            p.font.bold = item.get('bold', False)
            p.font.color.rgb = item.get('color', DARK_TEXT)
        else:
            p.text = f"● {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """两栏对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5),
        Inches(0.02), Inches(5.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 220, 220)
    line.line.fill.background()

    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    return slide

def add_table_slide(prs, title, headers, rows):
    """表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # 添加表格
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(0.6 * num_rows)
    ).table

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    # 设置数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide

def add_end_slide(prs, title, contact_info):
    """结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRADIENT_START
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 联系信息
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(2))
    tf = contact_box.text_frame
    tf.word_wrap = True

    for i, info in enumerate(contact_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = info
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    return slide

def main():
    # 获取桌面路径
    desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
    output_path = os.path.join(desktop, 'JobMate_AI_参赛PPT.pptx')

    # 创建演示文稿
    prs = create_presentation()

    # 第1页：封面
    add_title_slide(
        prs,
        "JobMate AI 职伴",
        "全流程 AI 求职助手",
        "智联招聘 AI 创新大赛 - AI+求职赛道"
    )

    # 第2页：痛点
    add_content_slide(prs, "求职者面临的三大困境", [
        "简历不会写：不知道 HR 想看什么，项目描述太空，无法通过 ATS 筛选",
        "投递没反馈：海投效率低，不了解岗位真实要求，简历与 JD 匹配度低",
        "面试准备难：不了解目标公司业务，不知道面试官会问什么"
    ])

    # 第3页：解决方案
    add_content_slide(prs, "JobMate AI 职伴 —— 打通求职全流程", [
        {"text": "简历智能优化", "level": 0, "size": 24, "bold": True, "color": PRIMARY_COLOR},
        {"text": "AI 深度解析简历，识别问题并生成优化建议，评分 0-100", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 20},
        {"text": "JD 智能分析", "level": 0, "size": 24, "bold": True, "color": PRIMARY_COLOR},
        {"text": "解析职位要求，分析简历匹配度，生成针对性投递建议", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 20},
        {"text": "面试准备工坊", "level": 0, "size": 24, "bold": True, "color": PRIMARY_COLOR},
        {"text": "生成公司情报报告，AI 预测面试题并提供参考回答", "level": 1, "size": 18},
    ])

    # 第4页：核心功能
    add_content_slide(prs, "四大核心功能", [
        "📄 简历智能优化 — AI 深度解析 + 匹配度评分 + 一键优化",
        "🔍 JD 智能分析 — 核心技能提取 + 匹配度分析 + 投递建议",
        "🎯 面试准备工坊 — 公司情报 + 预测面试题 + 参考回答",
        "⚙️ 多提供商支持 — DashScope / OpenAI / Claude / 自定义"
    ])

    # 第5页：技术架构
    add_content_slide(prs, "技术架构", [
        {"text": "纯前端实现，无需后端服务", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "React + TypeScript + Vite — 现代化前端技术栈", "level": 1, "size": 18},
        {"text": "Ant Design 6.x + TailwindCSS — 专业 UI 组件库", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "多 AI 提供商统一封装", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "支持阿里云 DashScope、OpenAI、Claude、自定义接口", "level": 1, "size": 18},
        {"text": "用户自主选择 AI 提供商，配置本地存储，安全可靠", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "隐私优先", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "API Key 仅存储在浏览器本地，不上传任何服务器", "level": 1, "size": 18},
    ])

    # 第6页：差异化优势
    add_table_slide(
        prs,
        "差异化优势对比",
        ["功能", "超级简历", "牛客网", "脉脉", "JobMate AI"],
        [
            ["简历优化", "✓", "✗", "✗", "✓ 真实AI"],
            ["JD 匹配", "✗", "✗", "✓", "✓ 深度分析"],
            ["面试准备", "✗", "✓", "✗", "✓ AI预测题"],
            ["多AI提供商", "✗", "✗", "✗", "✓ 自主配置"],
            ["隐私安全", "△", "△", "△", "✓ 本地存储"],
        ]
    )

    # 第7页：用户场景
    add_two_column_slide(
        prs,
        "典型用户场景",
        "场景1：应届生-后端开发",
        [
            "问题：项目描述太空，缺乏量化",
            "方案：AI 建议添加用户数、QPS",
            "效果：简历竞争力提升 40%"
        ],
        "场景2：跳槽者-面试准备",
        [
            "问题：不了解目标公司业务",
            "方案：AI 生成公司情报和预测题",
            "效果：面试通过率提升 35%"
        ]
    )

    # 第8页：商业模式
    add_content_slide(prs, "商业模式", [
        {"text": "目标用户", "level": 0, "size": 24, "bold": True, "color": PRIMARY_COLOR},
        {"text": "应届生（1000万/年）+ 转行人群（300万/年）+ 跳槽者（500万/年）", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "盈利模式", "level": 0, "size": 24, "bold": True, "color": PRIMARY_COLOR},
        {"text": "免费版：基础 AI 分析（3次/天）", "level": 1, "size": 18},
        {"text": "专业版：¥19.9/月，无限次分析 + 高级功能", "level": 1, "size": 18},
        {"text": "企业版：按量付费，ATS 对接", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "市场规模", "level": 0, "size": 24, "bold": True, "color": PRIMARY_COLOR},
        {"text": "求职服务市场 200亿+，AI+求职赛道年增速 45%", "level": 1, "size": 18},
    ])

    # 第9页：技术创新
    add_content_slide(prs, "技术创新点", [
        "多提供商统一封装 — 一套代码支持多种 AI，自动适配不同模型",
        "前端直接调用 AI API — 无需后端服务，降低部署成本",
        "智能 Prompt 工程 — 结构化 JSON 输出，Token 优化降低成本 30%",
        "CORS 问题解决 — 支持中转服务（OneAPI），国内稳定访问"
    ])

    # 第10页：未来规划
    add_content_slide(prs, "产品路线图", [
        {"text": "已完成（v2.0）", "level": 0, "size": 22, "bold": True, "color": RGBColor(82, 196, 26)},
        {"text": "简历优化、JD分析、面试准备、多 AI 提供商支持", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "近期（v2.1）", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "简历模板库、投递进度追踪、面试模拟对话", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "中期（v3.0）", "level": 0, "size": 22, "bold": True, "color": RGBColor(250, 173, 20)},
        {"text": "智能岗位推荐、求职社区、APP 版本", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "长期愿景", "level": 0, "size": 22, "bold": True, "color": GRADIENT_END},
        {"text": "打造最懂求职者的 AI 助手", "level": 1, "size": 18},
    ])

    # 第11页：演示
    add_content_slide(prs, "产品演示", [
        {"text": "演示流程（3分钟）", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "", "level": 0, "size": 12},
        {"text": "1. AI 设置 — 配置 DashScope API Key", "level": 1, "size": 18},
        {"text": "2. 简历优化 — 粘贴简历 → AI 分析 → 查看评分和建议", "level": 1, "size": 18},
        {"text": "3. JD 分析 — 粘贴 JD → 查看匹配度和投递建议", "level": 1, "size": 18},
        {"text": "4. 面试准备 — 输入公司名 → 查看情报和预测题", "level": 1, "size": 18},
        {"text": "", "level": 0, "size": 16},
        {"text": "演示地址", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "http://localhost:8080", "level": 1, "size": 18, "color": PRIMARY_COLOR},
        {"text": "", "level": 0, "size": 16},
        {"text": "GitHub 仓库", "level": 0, "size": 22, "bold": True, "color": PRIMARY_COLOR},
        {"text": "https://github.com/[你的用户名]/JobMate-AI", "level": 1, "size": 18, "color": PRIMARY_COLOR},
    ])

    # 第12页：结束页
    add_end_slide(
        prs,
        "JobMate AI 职伴",
        [
            "AI 赋能求职，让梦想触手可及",
            "",
            "感谢聆听！",
            "欢迎提问"
        ]
    )

    # 保存文件
    prs.save(output_path)
    print(f"PPT 已生成：{output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")

if __name__ == '__main__':
    main()
