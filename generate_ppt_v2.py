#!/usr/bin/env python3
"""
JobMate AI 参赛PPT - 优化版
更专业的布局和视觉设计
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# 配色方案 - 现代科技感
COLORS = {
    'primary': RGBColor(24, 144, 255),      # Ant Design Blue
    'primary_dark': RGBColor(0, 100, 200),
    'gradient_start': RGBColor(102, 126, 234),  # Purple
    'gradient_end': RGBColor(118, 75, 162),   # Dark Purple
    'accent': RGBColor(82, 196, 26),        # Success Green
    'warning': RGBColor(250, 173, 20),      # Warning Yellow
    'error': RGBColor(255, 77, 79),         # Error Red
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
    'dark': RGBColor(51, 51, 51),
    'gray': RGBColor(102, 102, 102),
    'light_gray': RGBColor(240, 240, 240),
    'bg_light': RGBColor(248, 250, 252),
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    return prs

def add_gradient_background(slide, prs):
    """添加渐变背景"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['gradient_start']
    shape.line.fill.background()
    # 将背景移到最底层
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_cover_slide(prs, title, subtitle, extra=""):
    """封面页 - 现代渐变设计"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 渐变背景
    add_gradient_background(slide, prs)

    # 装饰圆形
    for i, (x, y, size) in enumerate([(11, 5, 3), (10, 6, 2), (12, 4, 1.5)]):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['white']
        circle.fill.fore_color.brightness = 0.8
        circle.line.fill.background()
        # 透明度效果
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsmap

    # Logo 区域
    logo_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = COLORS['white']
    logo_box.line.fill.background()

    # Logo 文字
    logo_text = slide.shapes.add_textbox(
        Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    tf = logo_text.text_frame
    p = tf.paragraphs[0]
    p.text = "J"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['gradient_start']
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Inches(0.3)

    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.3), Inches(12.333), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 底部信息
    if extra:
        extra_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.6)
        )
        tf = extra_box.text_frame
        p = tf.paragraphs[0]
        p.text = extra
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, number, title):
    """章节分隔页 - 大号数字设计"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 左侧色块
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(0.8), prs.slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['primary']
    left_bar.line.fill.background()

    # 章节数字
    num_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(3), Inches(2)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{number}"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # 章节标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.8), Inches(10), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    return slide

def add_pain_points_slide(prs):
    """痛点页 - 三栏卡片布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "求职者面临的三大困境"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 三个痛点卡片
    cards = [
        ("简历不会写", ["不知道 HR 想看什么", "项目描述太空，缺乏量化", "无法通过 ATS 筛选"], COLORS['error']),
        ("投递没反馈", ["海投效率低，不知哪些岗位适合", "不了解岗位真实要求", "简历与 JD 匹配度低"], COLORS['warning']),
        ("面试准备难", ["不了解目标公司业务", "不知道面试官会问什么", "缺乏系统性准备方法"], COLORS['primary']),
    ]

    for i, (title, items, color) in enumerate(cards):
        x = 0.5 + i * 4.2

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3.8), Inches(5.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # 顶部色条
        color_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.6), Inches(3.8), Inches(0.1)
        )
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = color
        color_bar.line.fill.background()

        # 图标圆圈
        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 1.5), Inches(1.9), Inches(0.8), Inches(0.8)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = color
        icon_circle.line.fill.background()

        # 图标数字
        icon_text = slide.shapes.add_textbox(
            Inches(x + 1.5), Inches(1.95), Inches(0.8), Inches(0.7)
        )
        tf = icon_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 标题
        card_title = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(2.8), Inches(3.4), Inches(0.6)
        )
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # 内容列表
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(3.5), Inches(3.4), Inches(3.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['gray']
            p.space_after = Pt(12)

    return slide

def add_solution_slide(prs):
    """解决方案页 - 流程图设计"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "JobMate AI 职伴 —— 打通求职全流程"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 流程图背景
    flow_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5)
    )
    flow_bg.fill.solid()
    flow_bg.fill.fore_color.rgb = COLORS['bg_light']
    flow_bg.line.fill.background()

    # 三个主要模块
    modules = [
        ("📄 简历优化", "AI 深度解析\n识别问题\n生成建议", COLORS['primary']),
        ("🔍 JD 分析", "技能匹配\n投递建议\n竞争评估", COLORS['accent']),
        ("🎯 面试准备", "公司情报\n预测题目\n参考回答", COLORS['gradient_end']),
    ]

    for i, (title, desc, color) in enumerate(modules):
        x = 1.2 + i * 3.8

        # 模块卡片
        module = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(3.2), Inches(4.2)
        )
        module.fill.solid()
        module.fill.fore_color.rgb = COLORS['white']
        module.line.color.rgb = color
        module.line.width = Pt(3)

        # 顶部色块
        top_color = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(3.2), Inches(0.8)
        )
        top_color.fill.solid()
        top_color.fill.fore_color.rgb = color
        top_color.line.fill.background()

        # 标题
        module_title = slide.shapes.add_textbox(
            Inches(x), Inches(2.35), Inches(3.2), Inches(0.6)
        )
        tf = module_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 描述
        module_desc = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(3.2), Inches(2.8), Inches(3)
        )
        tf = module_desc.text_frame
        tf.word_wrap = True
        lines = desc.split('\n')
        for j, line in enumerate(lines):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {line}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['gray']
            p.space_after = Pt(16)

        # 箭头（除最后一个）
        if i < len(modules) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + 3.3), Inches(4), Inches(0.4), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['primary']
            arrow.line.fill.background()

    # 底部总结
    summary = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.5)
    )
    tf = summary.text_frame
    p = tf.paragraphs[0]
    p.text = "从简历到面试，每个环节都有 AI 辅助"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_features_slide(prs):
    """核心功能页 - 四宫格设计"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "四大核心功能"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 四个功能卡片
    features = [
        ("📄", "简历智能优化", "AI 深度解析 + 评分 0-100 + 一键优化", COLORS['primary']),
        ("🔍", "JD 智能分析", "技能匹配分析 + 投递建议 + 竞争力评估", COLORS['accent']),
        ("🎯", "面试准备工坊", "公司情报 + 预测面试题 + 参考回答", COLORS['warning']),
        ("⚙️", "多提供商支持", "DashScope / OpenAI / Claude / 自定义", COLORS['gradient_end']),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6
        y = 1.5 + row * 3

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.7), Inches(2.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # 图标区域
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.5), Inches(0.8), Inches(0.8)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = color
        icon_bg.line.fill.background()

        # 图标文字
        icon_text = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.55), Inches(0.8), Inches(0.7)
        )
        tf = icon_text.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

        # 标题
        feat_title = slide.shapes.add_textbox(
            Inches(x + 1.3), Inches(y + 0.5), Inches(4.2), Inches(0.6)
        )
        tf = feat_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color

        # 描述
        feat_desc = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 1.5), Inches(5.1), Inches(0.9)
        )
        tf = feat_desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['gray']

    return slide

def add_tech_arch_slide(prs):
    """技术架构页 - 分层架构图"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "技术架构"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 架构图容器
    arch_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(7), Inches(5.5)
    )
    arch_bg.fill.solid()
    arch_bg.fill.fore_color.rgb = COLORS['bg_light']
    arch_bg.line.fill.background()

    # 三层架构
    layers = [
        ("表现层", "React + TypeScript + Vite\nAnt Design + TailwindCSS", COLORS['primary']),
        ("服务层", "AI Service (ai.ts)\n多提供商统一封装", COLORS['accent']),
        ("数据层", "LocalStorage\nAPI Key 本地存储", COLORS['warning']),
    ]

    for i, (layer, desc, color) in enumerate(layers):
        y = 1.8 + i * 1.7

        # 层背景
        layer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(6.6), Inches(1.4)
        )
        layer_bg.fill.solid()
        layer_bg.fill.fore_color.rgb = COLORS['white']
        layer_bg.line.color.rgb = color
        layer_bg.line.width = Pt(2)

        # 层标题
        layer_title = slide.shapes.add_textbox(
            Inches(1.2), Inches(y + 0.15), Inches(2), Inches(0.5)
        )
        tf = layer_title.text_frame
        p = tf.paragraphs[0]
        p.text = layer
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # 层描述
        layer_desc = slide.shapes.add_textbox(
            Inches(1.2), Inches(y + 0.6), Inches(6.2), Inches(0.7)
        )
        tf = layer_desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['gray']

        # 箭头（除最后一个）
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, Inches(4), Inches(y + 1.45), Inches(0.5), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['gray']
            arrow.line.fill.background()

    # 右侧技术亮点
    highlights = [
        ("纯前端实现", "无需后端服务，部署简单", COLORS['primary']),
        ("隐私优先", "API Key 本地存储", COLORS['accent']),
        ("多提供商", "DashScope / OpenAI / Claude", COLORS['warning']),
        ("CORS 解决", "支持中转服务 OneAPI", COLORS['gradient_end']),
    ]

    for i, (title, desc, color) in enumerate(highlights):
        y = 1.6 + i * 1.4

        # 高亮卡片
        hl_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(y), Inches(4.8), Inches(1.2)
        )
        hl_card.fill.solid()
        hl_card.fill.fore_color.rgb = COLORS['white']
        hl_card.line.color.rgb = color
        hl_card.line.width = Pt(1.5)

        # 小色块
        color_block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8), Inches(y), Inches(0.15), Inches(1.2)
        )
        color_block.fill.solid()
        color_block.fill.fore_color.rgb = color
        color_block.line.fill.background()

        # 标题
        hl_title = slide.shapes.add_textbox(
            Inches(8.3), Inches(y + 0.15), Inches(4.3), Inches(0.4)
        )
        tf = hl_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # 描述
        hl_desc = slide.shapes.add_textbox(
            Inches(8.3), Inches(y + 0.55), Inches(4.3), Inches(0.5)
        )
        tf = hl_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['gray']

    return slide

def add_comparison_slide(prs):
    """竞品对比页 - 表格设计"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "差异化优势对比"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 表格数据
    headers = ["功能", "超级简历", "牛客网", "脉脉", "JobMate AI"]
    rows = [
        ["简历优化", "✓", "✗", "✗", "✓ 真实AI"],
        ["JD 匹配", "✗", "✗", "✓", "✓ 深度分析"],
        ["面试准备", "✗", "✓", "✗", "✓ AI预测题"],
        ["多AI提供商", "✗", "✗", "✗", "✓ 自主配置"],
        ["隐私安全", "△", "△", "△", "✓ 本地存储"],
    ]

    # 创建表格
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.5),
        Inches(11.7), Inches(0.7 * num_rows)
    ).table

    # 设置列宽
    col_widths = [2.5, 2.1, 2.1, 2.1, 2.9]
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)

            # 最后一列（我们的列）高亮
            if col_idx == num_cols - 1:
                p.font.bold = True
                p.font.color.rgb = COLORS['primary']
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['bg_light']
            else:
                p.font.color.rgb = COLORS['dark']

            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 核心优势
    advantages = [
        "全流程覆盖",
        "真实 AI 能力",
        "隐私优先",
        "灵活配置"
    ]

    adv_y = 6
    adv_title = slide.shapes.add_textbox(
        Inches(0.8), Inches(adv_y), Inches(2.5), Inches(0.5)
    )
    tf = adv_title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心优势："
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    for i, adv in enumerate(advantages):
        x = 3 + i * 2.5
        adv_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(adv_y - 0.05), Inches(2.3), Inches(0.6)
        )
        adv_box.fill.solid()
        adv_box.fill.fore_color.rgb = COLORS['primary']
        adv_box.line.fill.background()

        adv_text = slide.shapes.add_textbox(
            Inches(x), Inches(adv_y - 0.05), Inches(2.3), Inches(0.6)
        )
        tf = adv_text.text_frame
        p = tf.paragraphs[0]
        p.text = adv
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_business_slide(prs):
    """商业模式页 - 三栏卡片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "商业模式"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 目标用户
    user_title = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(3), Inches(0.5)
    )
    tf = user_title.text_frame
    p = tf.paragraphs[0]
    p.text = "目标用户"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    users = [
        ("应届生", "1000万/年", "缺乏经验"),
        ("转行人群", "300万/年", "经验不匹配"),
        ("跳槽者", "500万/年", "时间紧张"),
        ("海归", "50万/年", "不了解国内规则"),
    ]

    for i, (name, num, pain) in enumerate(users):
        y = 2 + i * 1.2

        # 用户卡片
        user_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(3.5), Inches(1)
        )
        user_card.fill.solid()
        user_card.fill.fore_color.rgb = COLORS['white']
        user_card.line.color.rgb = COLORS['light_gray']
        user_card.line.width = Pt(1)

        # 用户名
        name_text = slide.shapes.add_textbox(
            Inches(1), Inches(y + 0.1), Inches(1.5), Inches(0.4)
        )
        tf = name_text.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 数量
        num_text = slide.shapes.add_textbox(
            Inches(2.5), Inches(y + 0.1), Inches(1.5), Inches(0.4)
        )
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['primary']

        # 痛点
        pain_text = slide.shapes.add_textbox(
            Inches(1), Inches(y + 0.5), Inches(3.1), Inches(0.4)
        )
        tf = pain_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"痛点：{pain}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']

    # 盈利模式
    profit_title = slide.shapes.add_textbox(
        Inches(5), Inches(1.4), Inches(3), Inches(0.5)
    )
    tf = profit_title.text_frame
    p = tf.paragraphs[0]
    p.text = "盈利模式"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    profits = [
        ("免费版", "基础 AI 分析\n3次/天", COLORS['light_gray']),
        ("专业版", "¥19.9/月\n无限次 + 高级功能", COLORS['accent']),
        ("企业版", "按量付费\nATS 对接", COLORS['warning']),
    ]

    for i, (name, desc, color) in enumerate(profits):
        y = 2 + i * 1.5

        # 盈利卡片
        profit_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(y), Inches(3.5), Inches(1.3)
        )
        profit_card.fill.solid()
        profit_card.fill.fore_color.rgb = COLORS['white']
        profit_card.line.color.rgb = color
        profit_card.line.width = Pt(2)

        # 名称
        name_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(y + 0.1), Inches(1.2), Inches(0.4)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # 描述
        desc_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(y + 0.5), Inches(3.1), Inches(0.7)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc.replace('\n', ' | ')
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['gray']

    # 市场规模
    market_title = slide.shapes.add_textbox(
        Inches(9), Inches(1.4), Inches(3.5), Inches(0.5)
    )
    tf = market_title.text_frame
    p = tf.paragraphs[0]
    p.text = "市场规模"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['warning']

    # 市场规模数据
    market_data = [
        ("200亿+", "求职服务市场"),
        ("45%", "AI+求职赛道年增速"),
    ]

    for i, (num, label) in enumerate(market_data):
        y = 2.2 + i * 2

        # 数据展示
        num_box = slide.shapes.add_textbox(
            Inches(9), Inches(y), Inches(3.5), Inches(1)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['warning']

        label_box = slide.shapes.add_textbox(
            Inches(9), Inches(y + 0.9), Inches(3.5), Inches(0.5)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['gray']

    return slide

def add_demo_slide(prs):
    """演示页 - 截图占位"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "产品演示"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 三个截图占位框
    screenshots = [
        ("简历优化", "AI 分析 + 评分 + 建议"),
        ("JD 分析", "匹配度分析 + 投递建议"),
        ("面试准备", "公司情报 + 预测题"),
    ]

    for i, (title, desc) in enumerate(screenshots):
        x = 0.8 + i * 4.1

        # 截图框
        screenshot = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.8), Inches(4.5)
        )
        screenshot.fill.solid()
        screenshot.fill.fore_color.rgb = COLORS['bg_light']
        screenshot.line.color.rgb = COLORS['light_gray']
        screenshot.line.width = Pt(2)

        # 占位文字
        placeholder = slide.shapes.add_textbox(
            Inches(x), Inches(3.2), Inches(3.8), Inches(1)
        )
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"【{title}截图】"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # 标题
        screenshot_title = slide.shapes.add_textbox(
            Inches(x), Inches(6.1), Inches(3.8), Inches(0.5)
        )
        tf = screenshot_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

        # 描述
        screenshot_desc = slide.shapes.add_textbox(
            Inches(x), Inches(6.5), Inches(3.8), Inches(0.5)
        )
        tf = screenshot_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_end_slide(prs, title, subtitle):
    """结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 渐变背景
    add_gradient_background(slide, prs)

    # 装饰圆形
    for x, y, size in [(11, 5, 3), (10, 6, 2), (12, 4, 1.5), (-1, 6, 2.5), (0, 7, 1.5)]:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['white']
        circle.line.fill.background()

    # Logo
    logo_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = COLORS['white']
    logo_box.line.fill.background()

    logo_text = slide.shapes.add_textbox(
        Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    tf = logo_text.text_frame
    p = tf.paragraphs[0]
    p.text = "J"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['gradient_start']
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Inches(0.3)

    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 感谢文字
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6), Inches(12.333), Inches(0.8)
    )
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听！欢迎提问"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # 获取桌面路径
    desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
    output_path = os.path.join(desktop, 'JobMate_AI_参赛PPT_优化版.pptx')

    # 创建演示文稿
    prs = create_presentation()

    # 封面
    add_cover_slide(prs, "JobMate AI 职伴", "全流程 AI 求职助手", "智联招聘 AI 创新大赛 - AI+求职赛道")

    # 痛点
    add_pain_points_slide(prs)

    # 解决方案
    add_solution_slide(prs)

    # 核心功能
    add_features_slide(prs)

    # 技术架构
    add_tech_arch_slide(prs)

    # 竞品对比
    add_comparison_slide(prs)

    # 商业模式
    add_business_slide(prs)

    # 演示页（带截图占位）
    add_demo_slide(prs)

    # 结束页
    add_end_slide(prs, "JobMate AI 职伴", "AI 赋能求职，让梦想触手可及")

    # 保存文件
    prs.save(output_path)
    print(f"PPT 已生成：{output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")

if __name__ == '__main__':
    main()
